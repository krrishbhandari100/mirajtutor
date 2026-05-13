import io
import base64
import fitz
from docx import Document
from pptx import Presentation


def _image_to_base64(img_data: bytes) -> str:
    return base64.b64encode(img_data).decode("utf-8")


def parse_pdf(file_bytes: bytes):
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    text_parts = []
    pages = []

    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        text_parts.append(page.get_text())

        pix = page.get_pixmap(dpi=72)
        img_bytes = pix.tobytes("jpeg")
        pages.append({
            "number": page_num + 1,
            "image_base64": _image_to_base64(img_bytes),
        })

    return {
        "text": "\n\n".join(text_parts),
        "pages": pages,
        "total_pages": len(pages),
    }


def parse_pptx(file_bytes: bytes):
    prs = Presentation(io.BytesIO(file_bytes))
    text_parts = []
    images = []

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                slide_text.append(shape.text)
            if shape.shape_type == 13:
                try:
                    img = shape.image
                    img_bytes = img.blob
                    images.append({
                        "number": slide_num,
                        "image_base64": _image_to_base64(img_bytes),
                    })
                except Exception:
                    pass
        text_parts.append(f"--- Slide {slide_num} ---\n" + "\n".join(slide_text))

    return {
        "text": "\n\n".join(text_parts),
        "images": images,
        "total_slides": len(prs.slides),
    }


def parse_docx(file_bytes: bytes):
    doc = Document(io.BytesIO(file_bytes))
    text_parts = []
    images = []

    for para in doc.paragraphs:
        if para.text.strip():
            text_parts.append(para.text)

    for rel in doc.part.rels.values():
        if "image" in rel.reltype:
            try:
                img_bytes = rel.target_part.blob
                images.append({
                    "image_base64": _image_to_base64(img_bytes),
                })
            except Exception:
                pass

    return {
        "text": "\n\n".join(text_parts),
        "images": images,
    }


def parse_document(file_bytes: bytes, filename: str):
    ext = filename.lower().split(".")[-1] if "." in filename else ""

    if ext == "pdf":
        return {"type": "pdf", **parse_pdf(file_bytes)}
    elif ext in ("pptx", "ppt"):
        return {"type": "pptx", **parse_pptx(file_bytes)}
    elif ext in ("docx", "doc"):
        return {"type": "docx", **parse_docx(file_bytes)}
    else:
        raise ValueError(f"Unsupported file type: .{ext}")
